from __future__ import annotations

import csv
import io
import json
import os
import re
from pathlib import Path
from uuid import uuid4
from typing import Any
from urllib.error import HTTPError, URLError
from urllib.parse import parse_qs, urlparse
from urllib.request import Request, urlopen

try:
    from docx import Document
except Exception:  # pragma: no cover - optional dependency
    Document = None

try:
    from openpyxl import load_workbook
except Exception:  # pragma: no cover - optional dependency
    load_workbook = None

try:
    from pptx import Presentation
except Exception:  # pragma: no cover - optional dependency
    Presentation = None


TEXT_EXTENSIONS = {".csv", ".tsv", ".txt", ".md"}
SPREADSHEET_EXTENSIONS = {".xlsx"}
DOCUMENT_EXTENSIONS = {".docx"}
PRESENTATION_EXTENSIONS = {".pptx"}
GOOGLE_SHORTCUT_EXTENSIONS = {".gdoc", ".gsheet", ".gslides"}
SUPPORTED_EXTENSIONS = TEXT_EXTENSIONS | SPREADSHEET_EXTENSIONS | DOCUMENT_EXTENSIONS | PRESENTATION_EXTENSIONS | GOOGLE_SHORTCUT_EXTENSIONS

GOOGLE_DRIVE_EXPORTS = {
    ".gdoc": {
        "mime_type": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "suffix": ".docx",
    },
    ".gsheet": {
        "mime_type": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "suffix": ".xlsx",
    },
    ".gslides": {
        "mime_type": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "suffix": ".pptx",
    },
}

GOOGLE_DRIVE_ACCESS_TOKEN_ENV_VARS = (
    "FALCONBROOM_GOOGLE_DRIVE_ACCESS_TOKEN",
    "GOOGLE_DRIVE_ACCESS_TOKEN",
    "GOOGLE_OAUTH_ACCESS_TOKEN",
    "GOOGLE_OAUTH2_ACCESS_TOKEN",
)


def _safe_stem(name: str) -> str:
    base = Path(name).stem.strip() or "upload"
    cleaned = "".join(ch if ch.isalnum() or ch in ("-", "_") else "_" for ch in base)
    return cleaned or "upload"


def _decode_bytes(raw_bytes: bytes) -> str:
    for encoding in ("utf-8-sig", "utf-8", "cp1252", "latin-1"):
      try:
        return raw_bytes.decode(encoding)
      except Exception:
        continue
    return raw_bytes.decode("utf-8", errors="replace")


def _stringify(value: Any) -> str:
    if value is None:
        return ""
    if isinstance(value, bytes):
        return value.decode("utf-8", errors="replace")
    return str(value)


def _append_row(rows: list[dict[str, Any]], **row: Any) -> None:
    rows.append(
        {
            "source_kind": row.get("source_kind", "unknown"),
            "source_name": row.get("source_name", ""),
            "source_path": row.get("source_path", ""),
            "unit_kind": row.get("unit_kind", ""),
            "container_name": row.get("container_name", ""),
            "sheet_name": row.get("sheet_name", ""),
            "slide_number": row.get("slide_number", ""),
            "paragraph_index": row.get("paragraph_index", ""),
            "table_index": row.get("table_index", ""),
            "row_index": row.get("row_index", ""),
            "column_index": row.get("column_index", ""),
            "cell_label": row.get("cell_label", ""),
            "text": row.get("text", ""),
            "style_json": row.get("style_json", ""),
            "notes": row.get("notes", ""),
        }
    )


def _write_rows_to_csv(rows: list[dict[str, Any]], output_path: Path) -> None:
    fieldnames = [
        "source_kind",
        "source_name",
        "source_path",
        "unit_kind",
        "container_name",
        "sheet_name",
        "slide_number",
        "paragraph_index",
        "table_index",
        "row_index",
        "column_index",
        "cell_label",
        "text",
        "style_json",
        "notes",
    ]
    with output_path.open("w", newline="", encoding="utf-8") as handle:
        writer = csv.DictWriter(handle, fieldnames=fieldnames, quoting=csv.QUOTE_MINIMAL)
        writer.writeheader()
        for row in rows:
            writer.writerow(row)


def _sniff_delimiter(sample: str) -> str:
    if "\t" in sample and sample.count("\t") >= sample.count(","):
        return "\t"
    try:
        dialect = csv.Sniffer().sniff(sample, delimiters=",\t;|:")
        return dialect.delimiter
    except Exception:
        return "," if sample.count(",") >= sample.count(";") else ";"


def _first_nonempty(values: list[str | None]) -> str:
    for value in values:
        if value:
            return value
    return ""


def _google_drive_access_token() -> str:
    return _first_nonempty([os.environ.get(name, "").strip() for name in GOOGLE_DRIVE_ACCESS_TOKEN_ENV_VARS])


def _google_drive_file_id(candidate: str) -> str:
    if not candidate:
        return ""

    text = candidate.strip()
    parsed = urlparse(text)
    if parsed.query:
        query_params = parse_qs(parsed.query)
        for key in ("id", "fileId"):
            values = query_params.get(key)
            if values and values[0]:
                return values[0]

    patterns = [
        r"/d/([A-Za-z0-9_-]{8,})",
        r"[?&]id=([A-Za-z0-9_-]{8,})",
        r"/folders/([A-Za-z0-9_-]{8,})",
        r"^([A-Za-z0-9_-]{8,})$",
    ]
    for pattern in patterns:
        match = re.search(pattern, text)
        if match:
            return match.group(1)
    return ""


def _google_shortcut_payload(source_path: Path) -> dict[str, Any]:
    try:
        payload = json.loads(_decode_bytes(source_path.read_bytes()))
        return payload if isinstance(payload, dict) else {}
    except Exception:
        return {}


def _download_google_drive_export(file_id: str, mime_type: str) -> bytes:
    token = _google_drive_access_token()
    if not token:
        raise RuntimeError(
            "Google Drive access token not configured. Set FALCONBROOM_GOOGLE_DRIVE_ACCESS_TOKEN to enable authenticated fetches."
        )

    export_url = f"https://www.googleapis.com/drive/v3/files/{file_id}/export?mimeType={mime_type}"
    request = Request(
        export_url,
        headers={
            "Authorization": f"Bearer {token}",
            "Accept": "application/octet-stream",
        },
    )
    try:
        with urlopen(request, timeout=60) as response:
            return response.read()
    except HTTPError as exc:
        detail = exc.read().decode("utf-8", errors="replace") if exc.fp else str(exc)
        raise RuntimeError(f"Drive export failed: {detail or exc}") from exc
    except URLError as exc:
        raise RuntimeError(f"Drive export request failed: {exc}") from exc


def _extract_google_drive_export(source_path: Path, output_dir: Path, rows: list[dict[str, Any]], warnings: list[str]) -> bool:
    payload = _google_shortcut_payload(source_path)
    url = _first_nonempty(
        [
            _stringify(payload.get("url")),
            _stringify(payload.get("doc_url")),
            _stringify(payload.get("sheet_url")),
            _stringify(payload.get("presentation_url")),
        ]
    )
    resource_id = _first_nonempty(
        [
            _stringify(payload.get("resource_id")),
            _stringify(payload.get("doc_id")),
            _stringify(payload.get("sheet_id")),
            _stringify(payload.get("presentation_id")),
            _stringify(payload.get("id")),
        ]
    )
    file_id = _google_drive_file_id(resource_id or url)
    export_info = GOOGLE_DRIVE_EXPORTS.get(source_path.suffix.lower())

    if not export_info or not file_id:
        return False

    try:
        export_bytes = _download_google_drive_export(file_id, export_info["mime_type"])
    except Exception as exc:
        warnings.append(str(exc))
        return False

    exported_path = output_dir / f"{_safe_stem(source_path.name)}_{uuid4().hex[:8]}{export_info['suffix']}"
    exported_path.write_bytes(export_bytes)
    warnings.append(f"Fetched Google Drive export for {source_path.name}.")

    if export_info["suffix"] == ".docx":
        _extract_docx_rows(exported_path, rows)
    elif export_info["suffix"] == ".xlsx":
        _extract_workbook_rows(exported_path, rows)
    elif export_info["suffix"] == ".pptx":
        _extract_pptx_rows(exported_path, rows)
    else:
        _extract_text_rows(exported_path, "google_drive", rows)

    return True


def _extract_text_rows(source_path: Path, source_kind: str, rows: list[dict[str, Any]]) -> None:
    raw_text = _decode_bytes(source_path.read_bytes())
    stripped = raw_text.strip()
    if not stripped:
        _append_row(
            rows,
            source_kind=source_kind,
            source_name=source_path.name,
            source_path=str(source_path),
            unit_kind="line",
            row_index=1,
            text="",
            notes="empty text file",
        )
        return

    sample = "\n".join(line for line in raw_text.splitlines() if line.strip())[:4096]
    delimiter = _sniff_delimiter(sample) if sample else ","
    reader = csv.reader(io.StringIO(raw_text), delimiter=delimiter, quotechar='"', escapechar="\\", strict=False)

    any_rows = False
    for row_index, row in enumerate(reader, start=1):
        any_rows = True
        if not row:
            continue
        row_text = delimiter.join(cell for cell in row if cell is not None)
        _append_row(
            rows,
            source_kind=source_kind,
            source_name=source_path.name,
            source_path=str(source_path),
            unit_kind="line",
            row_index=row_index,
            text=row_text,
            style_json=json.dumps({"delimiter": delimiter, "fields": len(row)}, ensure_ascii=False),
        )

    if not any_rows:
        for row_index, line in enumerate(raw_text.splitlines(), start=1):
            if not line.strip():
                continue
            _append_row(
                rows,
                source_kind=source_kind,
                source_name=source_path.name,
                source_path=str(source_path),
                unit_kind="line",
                row_index=row_index,
                text=line,
                notes="line fallback",
            )


def _extract_workbook_rows(source_path: Path, rows: list[dict[str, Any]]) -> None:
    if load_workbook is None:
        raise RuntimeError("openpyxl is not installed. Install requirements to enable spreadsheet extraction.")

    workbook = load_workbook(source_path, data_only=False, read_only=True)
    for sheet_name in workbook.sheetnames:
        sheet = workbook[sheet_name]
        for row_index, row in enumerate(sheet.iter_rows(), start=1):
            for column_index, cell in enumerate(row, start=1):
                try:
                    value = cell.value
                    if value is None or value == "":
                        continue

                    style = {
                        "data_type": cell.data_type,
                        "bold": bool(getattr(cell.font, "bold", False)),
                        "italic": bool(getattr(cell.font, "italic", False)),
                        "font_name": getattr(cell.font, "name", None),
                        "font_size": getattr(cell.font, "sz", None),
                        "number_format": getattr(cell, "number_format", None),
                        "alignment": getattr(cell.alignment, "horizontal", None),
                        "fill": getattr(getattr(cell.fill, "fgColor", None), "rgb", None),
                    }
                    if isinstance(value, str) and value.startswith("="):
                        style["formula"] = value

                    _append_row(
                        rows,
                        source_kind="xlsx",
                        source_name=source_path.name,
                        source_path=str(source_path),
                        unit_kind="cell",
                        sheet_name=sheet_name,
                        row_index=row_index,
                        column_index=column_index,
                        cell_label=cell.coordinate,
                        text=_stringify(value),
                        style_json=json.dumps(style, ensure_ascii=False, default=_stringify),
                    )
                except Exception as exc:
                    _append_row(
                        rows,
                        source_kind="xlsx",
                        source_name=source_path.name,
                        source_path=str(source_path),
                        unit_kind="cell_error",
                        sheet_name=sheet_name,
                        row_index=row_index,
                        column_index=column_index,
                        cell_label=getattr(cell, "coordinate", ""),
                        text=_stringify(getattr(cell, "value", "")),
                        notes=str(exc),
                    )


def _extract_docx_rows(source_path: Path, rows: list[dict[str, Any]]) -> None:
    if Document is None:
        raise RuntimeError("python-docx is not installed. Install requirements to enable Word extraction.")

    document = Document(source_path)
    for paragraph_index, paragraph in enumerate(document.paragraphs, start=1):
        text = paragraph.text.strip()
        if not text:
            continue
        run_bold_count = sum(1 for run in paragraph.runs if getattr(run, "bold", False))
        run_italic_count = sum(1 for run in paragraph.runs if getattr(run, "italic", False))
        style = {
            "style_name": getattr(getattr(paragraph, "style", None), "name", None),
            "alignment": getattr(paragraph.paragraph_format.alignment, "value", None) if paragraph.paragraph_format else None,
            "run_count": len(paragraph.runs),
            "bold_runs": run_bold_count,
            "italic_runs": run_italic_count,
        }
        _append_row(
            rows,
            source_kind="docx",
            source_name=source_path.name,
            source_path=str(source_path),
            unit_kind="paragraph",
            paragraph_index=paragraph_index,
            text=text,
            style_json=json.dumps(style, ensure_ascii=False, default=_stringify),
        )

    for table_index, table in enumerate(document.tables, start=1):
        for row_index, table_row in enumerate(table.rows, start=1):
            for column_index, cell in enumerate(table_row.cells, start=1):
                try:
                    text = cell.text.strip()
                    if not text:
                        continue
                    style = {
                        "paragraphs": len(cell.paragraphs),
                        "tables": len(cell.tables),
                    }
                    _append_row(
                        rows,
                        source_kind="docx",
                        source_name=source_path.name,
                        source_path=str(source_path),
                        unit_kind="table_cell",
                        container_name=f"table_{table_index}",
                        table_index=table_index,
                        row_index=row_index,
                        column_index=column_index,
                        cell_label=f"r{row_index}c{column_index}",
                        text=text,
                        style_json=json.dumps(style, ensure_ascii=False, default=_stringify),
                    )
                except Exception as exc:
                    _append_row(
                        rows,
                        source_kind="docx",
                        source_name=source_path.name,
                        source_path=str(source_path),
                        unit_kind="table_cell_error",
                        container_name=f"table_{table_index}",
                        table_index=table_index,
                        row_index=row_index,
                        column_index=column_index,
                        cell_label=f"r{row_index}c{column_index}",
                        notes=str(exc),
                    )


def _extract_pptx_rows(source_path: Path, rows: list[dict[str, Any]]) -> None:
    if Presentation is None:
        raise RuntimeError("python-pptx is not installed. Install requirements to enable PowerPoint extraction.")

    presentation = Presentation(source_path)
    for slide_number, slide in enumerate(presentation.slides, start=1):
        for shape_index, shape in enumerate(slide.shapes, start=1):
            try:
                shape_name = getattr(shape, "name", f"shape_{shape_index}")
                if getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip():
                    for paragraph_index, paragraph in enumerate(shape.text_frame.paragraphs, start=1):
                        text = paragraph.text.strip()
                        if not text:
                            continue
                        style = {
                            "shape_name": shape_name,
                            "shape_type": getattr(shape, "shape_type", None),
                            "paragraph_level": getattr(paragraph, "level", None),
                            "alignment": getattr(paragraph.alignment, "value", None),
                            "run_count": len(paragraph.runs),
                        }
                        _append_row(
                            rows,
                            source_kind="pptx",
                            source_name=source_path.name,
                            source_path=str(source_path),
                            unit_kind="paragraph",
                            slide_number=slide_number,
                            container_name=shape_name,
                            paragraph_index=paragraph_index,
                            text=text,
                            style_json=json.dumps(style, ensure_ascii=False, default=_stringify),
                        )

                if getattr(shape, "has_table", False):
                    table = shape.table
                    for row_index, table_row in enumerate(table.rows, start=1):
                        for column_index, cell in enumerate(table_row.cells, start=1):
                            text = cell.text.strip()
                            if not text:
                                continue
                            _append_row(
                                rows,
                                source_kind="pptx",
                                source_name=source_path.name,
                                source_path=str(source_path),
                                unit_kind="table_cell",
                                slide_number=slide_number,
                                container_name=shape_name,
                                row_index=row_index,
                                column_index=column_index,
                                cell_label=f"r{row_index}c{column_index}",
                                text=text,
                                notes="slide table",
                            )
            except Exception as exc:
                _append_row(
                    rows,
                    source_kind="pptx",
                    source_name=source_path.name,
                    source_path=str(source_path),
                    unit_kind="shape_error",
                    slide_number=slide_number,
                    container_name=getattr(shape, "name", f"shape_{shape_index}"),
                    notes=str(exc),
                )


def _extract_google_shortcut_rows(source_path: Path, output_dir: Path, rows: list[dict[str, Any]], warnings: list[str]) -> None:
    payload = _google_shortcut_payload(source_path)

    if isinstance(payload, dict):
        url = payload.get("url") or payload.get("doc_url") or payload.get("sheet_url") or payload.get("presentation_url")
        title = payload.get("title") or payload.get("name") or source_path.stem
        resource_id = payload.get("resource_id") or payload.get("doc_id") or payload.get("id") or ""
        _append_row(
            rows,
            source_kind="google_shortcut",
            source_name=source_path.name,
            source_path=str(source_path),
            unit_kind="shortcut_metadata",
            text=_stringify(title),
            style_json=json.dumps(payload, ensure_ascii=False, default=_stringify),
            notes=f"{url or 'Google Drive shortcut'}; content extraction will use Drive export when authenticated",
        )
        if resource_id:
            _append_row(
                rows,
                source_kind="google_shortcut",
                source_name=source_path.name,
                source_path=str(source_path),
                unit_kind="shortcut_reference",
                text=_stringify(resource_id),
                notes=url or "",
            )
        if _extract_google_drive_export(source_path, output_dir, rows, warnings):
            return
        return

    _extract_text_rows(source_path, "google_shortcut", rows)


def convert_uploaded_file(source_path: str | Path, output_dir: str | Path) -> dict[str, Any]:
    source_path = Path(source_path)
    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)

    suffix = source_path.suffix.lower()
    target = output_dir / f"{_safe_stem(source_path.name)}_{source_path.stat().st_size}_{source_path.stat().st_mtime_ns & 0xFFFF:x}.csv"
    rows: list[dict[str, Any]] = []
    warnings: list[str] = []

    if suffix in TEXT_EXTENSIONS:
        _extract_text_rows(source_path, suffix.lstrip("."), rows)
    elif suffix in SPREADSHEET_EXTENSIONS:
        _extract_workbook_rows(source_path, rows)
    elif suffix in DOCUMENT_EXTENSIONS:
        _extract_docx_rows(source_path, rows)
    elif suffix in PRESENTATION_EXTENSIONS:
        _extract_pptx_rows(source_path, rows)
    elif suffix in GOOGLE_SHORTCUT_EXTENSIONS:
        _extract_google_shortcut_rows(source_path, output_dir, rows, warnings)
    else:
        warnings.append(f"Unsupported extension {suffix or '(none)'}; stored as line-based text fallback.")
        _extract_text_rows(source_path, suffix.lstrip(".") or "binary", rows)

    if not rows:
        warnings.append("No structured content was extracted; added a placeholder row.")
        _append_row(
            rows,
            source_kind=suffix.lstrip(".") or "unknown",
            source_name=source_path.name,
            source_path=str(source_path),
            unit_kind="placeholder",
            text="",
            notes="no structured content extracted",
        )

    _write_rows_to_csv(rows, target)

    metadata = {
        "source_path": str(source_path.resolve()),
        "normalized_path": str(target.resolve()),
        "source_kind": suffix.lstrip(".") or "unknown",
        "source_name": source_path.name,
        "row_count": len(rows),
        "warnings": warnings,
    }
    metadata_path = target.with_suffix(".json")
    metadata_path.write_text(json.dumps(metadata, indent=2, ensure_ascii=False), encoding="utf-8")
    metadata["metadata_path"] = str(metadata_path.resolve())
    metadata["path"] = metadata["normalized_path"]
    return metadata